# presentation/create_presentation.py
from pptx import Presentation

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Capstone: Business Analysis"
slide.placeholders[1].text = "E-Commerce, Real Estate & Churn — Executive Summary"
prs.slides.add_slide(prs.slide_layouts[1]).shapes.title.text = "Key Findings"
prs.save("presentation/business_presentation.pptx")
print("Saved presentation/business_presentation.pptx")
